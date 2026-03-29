from pptx import Presentation
from pptx.util import Pt


class PptxProcessor:
    """PowerPoint (.pptx) 文档处理器"""

    def parse(self, file_path: str) -> dict:
        """解析 pptx 文档结构"""
        prs = Presentation(file_path)
        slides = []
        for si, slide in enumerate(prs.slides):
            shapes = []
            for shi, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                paragraphs = []
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    paragraphs.append(
                        {
                            "index": pi,
                            "text": para.text,
                            "level": para.level,
                            "runs": [
                                {
                                    "text": run.text,
                                    "bold": run.font.bold,
                                    "italic": run.font.italic,
                                    "font_name": run.font.name,
                                    "font_size": str(run.font.size) if run.font.size else None,
                                }
                                for run in para.runs
                            ],
                        }
                    )
                shapes.append(
                    {
                        "shape_index": shi,
                        "shape_name": shape.name,
                        "text": shape.text_frame.text,
                        "paragraphs": paragraphs,
                    }
                )
            slides.append({"index": si, "shapes": shapes})

        return {
            "file_path": file_path,
            "format": "pptx",
            "slide_count": len(prs.slides),
            "slides": slides,
        }

    def extract_text(self, file_path: str) -> str:
        """提取纯文本"""
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)

    def get_stats(self, file_path: str) -> dict:
        """文档统计"""
        prs = Presentation(file_path)
        shape_count = 0
        total_chars = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_count += 1
                    total_chars += len(shape.text_frame.text)
        return {
            "slide_count": len(prs.slides),
            "shape_count": shape_count,
            "total_chars": total_chars,
        }

    def replace_text(self, file_path: str, replacements: dict, output_path: str) -> str:
        """替换指定位置的文本，保留格式

        replacements 格式: {"slide_idx:shape_idx:para_idx": "new_text"}
        shape_idx 是 has_text_frame 的 shape 中的顺序索引
        """
        prs = Presentation(file_path)

        # 建立 text shape 索引映射
        slide_shape_map: dict[int, list] = {}
        for si, slide in enumerate(prs.slides):
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            slide_shape_map[si] = text_shapes

        for key, new_text in replacements.items():
            parts = str(key).split(":")
            if len(parts) == 3:
                si, shi, pi = int(parts[0]), int(parts[1]), int(parts[2])
            elif len(parts) == 2:
                # slide_idx:para_idx — 在所有 shape 的段落中按全局索引
                si, pi = int(parts[0]), int(parts[1])
                shi = None
            else:
                continue

            text_shapes = slide_shape_map.get(si, [])
            if shi is not None:
                if 0 <= shi < len(text_shapes):
                    shape = text_shapes[shi]
                    paras = shape.text_frame.paragraphs
                    if 0 <= pi < len(paras):
                        self._replace_paragraph(paras[pi], new_text)
            else:
                # 按全局段落索引查找
                global_idx = 0
                for shape in text_shapes:
                    for para in shape.text_frame.paragraphs:
                        if global_idx == pi:
                            self._replace_paragraph(para, new_text)
                        global_idx += 1

        prs.save(output_path)
        return output_path

    def apply_format_changes(self, file_path: str, changes: list[dict], output_path: str) -> str:
        """应用格式修改指令

        changes 格式:
        [
            {"target": "font_size", "scope": "all", "value_pt": 14},
            {"target": "font_name", "scope": [0, 1], "value": "微软雅黑"},
            {"target": "line_spacing", "scope": "all", "value": 1.15},
        ]
        scope: "all" 或 slide 索引列表
        """
        prs = Presentation(file_path)
        all_slides = list(prs.slides)

        for change in changes:
            target = change.get("target")
            scope = change.get("scope", "all")
            slides = self._resolve_scope(all_slides, scope)

            for slide in slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        if target == "font_size":
                            value = Pt(change["value_pt"])
                            for run in para.runs:
                                run.font.size = value
                        elif target == "font_name":
                            name = change["value"]
                            for run in para.runs:
                                run.font.name = name
                        elif target == "line_spacing":
                            para.line_spacing = change["value"]
                        elif target == "space_after":
                            para.space_after = Pt(change["value_pt"])

        prs.save(output_path)
        return output_path

    def _replace_paragraph(self, para, new_text: str):
        """替换段落文本，保留第一个 run 的格式"""
        if para.runs:
            first_run = para.runs[0]
            bold = first_run.font.bold
            italic = first_run.font.italic
            font_name = first_run.font.name
            font_size = first_run.font.size
            for run in para.runs:
                run.text = ""
            first_run.text = new_text
            first_run.font.bold = bold
            first_run.font.italic = italic
            if font_name:
                first_run.font.name = font_name
            if font_size:
                first_run.font.size = font_size
        else:
            para.text = new_text

    def _resolve_scope(self, slides, scope):
        """解析 scope: "all" 或索引列表"""
        if scope == "all":
            return list(slides)
        if isinstance(scope, list):
            return [slides[i] for i in scope if 0 <= i < len(slides)]
        return []
