"""Agent 角色定义 + 自定义工具"""

import json
import logging
from pathlib import Path

from claude_agent_sdk import AgentDefinition, create_sdk_mcp_server, tool

from backend.processors import get_processor

logger = logging.getLogger("docflow.agents")


def _error_result(msg: str) -> dict:
    return {"content": [{"type": "text", "text": f"Error: {msg}"}]}


# === Agent 角色 ===

CONTENT_GENERATOR = AgentDefinition(
    description="内容生成器：根据用户需求生成文档初稿",
    prompt="""你是内容生成器。你的职责是根据用户的需求描述生成结构化的文档初稿。

规则:
1. 使用 write_document 工具将生成的内容写入指定路径
2. 生成的内容应结构清晰，包含标题、段落、列表等
3. 使用中文写作，符合学术/办公文档规范
4. 生成完成后，发消息给 content-editor 通知文档路径

不要生成 AI 味过重的内容（避免"此外"、"值得注意的是"等套话），但这不是你的核心职责，
后续编辑会进一步去味。
""",
    tools=["Read", "Write", "mcp__docflow-tools__write_document"],
    model="sonnet",
)

CONTENT_EDITOR = AgentDefinition(
    description="内容编辑：去 AI 味 + 注入文风 DNA",
    prompt="""你是内容编辑。你的职责是消除文档中的 AI 味道。

去 AI 味 Checklist:
- "此外"、"值得注意的是"、"综上所述" → 删除或用具体连接替代
- "在...方面"、"对于...而言" → 直接说主语+动词
- 被动句式（"被认为是"） → 转主动句
- 过度对仗的排比句 → 打破对称，长短句混搭
- 每段首句都是概括句 → 变化开头方式
- "首先...其次...最后..." → 减少明显的枚举标记

工作流程:
1. 使用 parse_document 工具读取文档结构
2. 逐段改写，去除 AI 味
3. 使用 replace_content 工具写入修改
4. 发消息给 format-designer 通知修改完成

收到 quality-reviewer 的返工消息时，只修改审核员标记的段落，不要重做全文。
""",
    tools=["Read", "mcp__docflow-tools__parse_document", "mcp__docflow-tools__replace_content"],
    model="sonnet",
)

FORMAT_DESIGNER = AgentDefinition(
    description="格式设计师：输出 JSON 格式修改指令，人类化排版",
    prompt="""你是格式设计师。你的职责是让文档排版更像人类手动排版的效果。

格式人类化原理: AI 生成的文档排版过于完美均匀，人类手动排版天然有微小不一致。
你的任务是引入自然偏差，而非制造混乱。

格式操作 (适用于所有文档类型):
- 字体大小: 在 ±1pt 范围内微调
- 字体名称: 选择合适的中文字体
- 行距/段间距: ±0.05 倍随机偏移

工作流程:
1. 使用 parse_document 分析当前排版
2. 生成 JSON 格式修改指令
3. 使用 apply_format 工具确定性执行
4. 发消息给 quality-reviewer 通知格式修改完成

所有变更限制在不破坏可读性的范围内。
""",
    tools=["Read", "mcp__docflow-tools__parse_document", "mcp__docflow-tools__apply_format"],
    model="sonnet",
)

QUALITY_REVIEWER = AgentDefinition(
    description="质量审核员：盲审文档质量，按维度打分",
    prompt="""你是质量审核员。你的职责是盲审文档质量。

评分维度 (总分 = 加权平均):
| 维度 | 权重 |
|------|------|
| 词汇自然度 | 30% |
| 句式多样性 | 20% |
| 格式人类感 | 25% |
| 逻辑连贯性 | 15% |
| 领域适配度 | 10% |

通过标准: 总分 >= 8.0

工作流程:
1. 使用 parse_document 读取精修后的文档
2. 逐维度评分 (1-10)
3. 使用 submit_score 工具提交评分结果

如果总分 < 8.0:
- 标记具体问题位置（段落索引 + 问题描述）
- 发消息给 content-editor: "返工要求: [具体段落和问题]"
- 如果是格式问题，也发消息给 format-designer

如果总分 >= 8.0:
- 发消息报告通过

最多审核 3 轮。第 3 轮后无论分数如何，输出当前版本。
""",
    tools=["Read", "mcp__docflow-tools__parse_document", "mcp__docflow-tools__submit_score"],
    model="sonnet",
)

AGENTS = {
    "content-generator": CONTENT_GENERATOR,
    "content-editor": CONTENT_EDITOR,
    "format-designer": FORMAT_DESIGNER,
    "quality-reviewer": QUALITY_REVIEWER,
}


# === 自定义工具 ===


@tool("parse_document", "解析文档结构（支持 docx/pptx/xlsx/pdf）", {"file_path": str})
async def parse_document_tool(args):
    try:
        file_path = args.get("file_path", "")
        if not file_path or not Path(file_path).exists():
            return _error_result(f"文件不存在: {file_path}")
        processor = get_processor(file_path)
        result = processor.parse(file_path)
        return {"content": [{"type": "text", "text": json.dumps(result, ensure_ascii=False)}]}
    except Exception as e:
        logger.error("parse_document failed: %s", e, exc_info=True)
        return _error_result(f"{type(e).__name__}: {e}")


@tool(
    "replace_content",
    "替换文档文本内容（保留格式，支持 docx/pptx/xlsx/pdf）",
    {"file_path": str, "output_path": str, "replacements": str},
)
async def replace_content_tool(args):
    try:
        file_path = args.get("file_path", "")
        if not file_path or not Path(file_path).exists():
            return _error_result(f"文件不存在: {file_path}")
        replacements = json.loads(args.get("replacements", "{}"))
        processor = get_processor(file_path)
        out = processor.replace_text(file_path, replacements, args.get("output_path", ""))
        return {"content": [{"type": "text", "text": f"内容替换完成，输出: {out}"}]}
    except json.JSONDecodeError as e:
        return _error_result(f"replacements JSON 格式错误: {e}")
    except Exception as e:
        logger.error("replace_content failed: %s", e, exc_info=True)
        return _error_result(f"{type(e).__name__}: {e}")


@tool(
    "apply_format",
    "执行 JSON 格式修改指令（支持 docx/pptx/xlsx/pdf）",
    {"file_path": str, "output_path": str, "changes": str},
)
async def apply_format_tool(args):
    try:
        file_path = args.get("file_path", "")
        if not file_path or not Path(file_path).exists():
            return _error_result(f"文件不存在: {file_path}")
        changes = json.loads(args.get("changes", "[]"))
        processor = get_processor(file_path)
        out = processor.apply_format_changes(file_path, changes, args.get("output_path", ""))
        return {"content": [{"type": "text", "text": f"格式修改已应用，输出: {out}"}]}
    except json.JSONDecodeError as e:
        return _error_result(f"changes JSON 格式错误: {e}")
    except Exception as e:
        logger.error("apply_format failed: %s", e, exc_info=True)
        return _error_result(f"{type(e).__name__}: {e}")


@tool("write_document", "将文本内容写入文档文件（支持 docx/pptx/xlsx）", {"output_path": str, "content": str})
async def write_document_tool(args):
    try:
        output_path = args.get("output_path", "")
        content = args.get("content", "")
        if not output_path:
            return _error_result("output_path is required")
        if not content:
            return _error_result("content is required")
        ext = Path(output_path).suffix.lower()

        if ext == ".pptx":
            _write_pptx(output_path, content)
        elif ext == ".xlsx":
            _write_xlsx(output_path, content)
        else:
            _write_docx(output_path, content)

        return {"content": [{"type": "text", "text": f"文档已创建: {output_path}"}]}
    except Exception as e:
        logger.error("write_document failed: %s", e, exc_info=True)
        return _error_result(f"{type(e).__name__}: {e}")


def _write_docx(output_path: str, content: str):
    from docx import Document

    doc = Document()
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(output_path)


def _write_pptx(output_path: str, content: str):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    current_tf = None

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# ") or line.startswith("## "):
            # 新页
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
            current_tf = tx_box.text_frame
            current_tf.word_wrap = True
            current_tf.text = line.lstrip("#").strip()
        else:
            if current_tf is None:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
                current_tf = tx_box.text_frame
                current_tf.word_wrap = True
                current_tf.text = line
            else:
                p = current_tf.add_paragraph()
                p.text = line.lstrip("- ")

    if not prs.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(output_path)


def _write_xlsx(output_path: str, content: str):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    row = 1
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# ") or line.startswith("## "):
            ws.cell(row=row, column=1, value=line.lstrip("#").strip())
        elif "|" in line and not line.startswith("|-"):
            # 表格行
            cells = [c.strip() for c in line.strip("|").split("|")]
            for ci, val in enumerate(cells, 1):
                ws.cell(row=row, column=ci, value=val)
        else:
            ws.cell(row=row, column=1, value=line.lstrip("- "))
        row += 1
    wb.save(output_path)
    wb.close()


@tool("submit_score", "提交质量评分结果", {"scores": str})
async def submit_score_tool(args):
    try:
        scores = json.loads(args.get("scores", "{}"))
        weights = {
            "vocabulary_naturalness": 0.3,
            "sentence_diversity": 0.2,
            "format_humanity": 0.25,
            "logical_coherence": 0.15,
            "domain_adaptation": 0.1,
        }
        total = sum(scores.get(k, 0) * w for k, w in weights.items())
        passed = total >= 8.0
        result = {**scores, "total": round(total, 1), "passed": passed}
        return {"content": [{"type": "text", "text": json.dumps(result, ensure_ascii=False)}]}
    except json.JSONDecodeError as e:
        return _error_result(f"scores JSON 格式错误: {e}")
    except Exception as e:
        logger.error("submit_score failed: %s", e, exc_info=True)
        return _error_result(f"{type(e).__name__}: {e}")


# MCP Server
docflow_tools = create_sdk_mcp_server(
    name="docflow-tools",
    version="1.0.0",
    tools=[parse_document_tool, replace_content_tool, apply_format_tool, write_document_tool, submit_score_tool],
)
