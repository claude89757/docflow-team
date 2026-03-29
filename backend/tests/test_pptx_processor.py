"""PptxProcessor 单元测试"""

import os
import tempfile

from pptx import Presentation
from pptx.util import Pt

from backend.processors.pptx_processor import PptxProcessor


def _create_simple_pptx(path: str, slide_texts: list[list[str]]):
    """创建测试用 pptx，每个 slide 包含一个文本框，每个字符串一个段落"""
    prs = Presentation()
    for texts in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        from pptx.util import Inches

        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = tx_box.text_frame
        tf.text = texts[0]
        for text in texts[1:]:
            p = tf.add_paragraph()
            p.text = text
    prs.save(path)


def test_parse():
    processor = PptxProcessor()
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        _create_simple_pptx(path, [["标题", "第一段内容"], ["第二页", "更多内容"]])

        result = processor.parse(path)
        assert result["format"] == "pptx"
        assert result["slide_count"] == 2
        assert len(result["slides"]) == 2
        assert len(result["slides"][0]["shapes"]) >= 1
        assert result["slides"][0]["shapes"][0]["paragraphs"][0]["text"] == "标题"


def test_extract_text():
    processor = PptxProcessor()
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        _create_simple_pptx(path, [["测试文本", "更多文本"]])

        text = processor.extract_text(path)
        assert "测试文本" in text
        assert "更多文本" in text


def test_get_stats():
    processor = PptxProcessor()
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        _create_simple_pptx(path, [["段落一", "段落二"], ["段落三"]])

        stats = processor.get_stats(path)
        assert stats["slide_count"] == 2
        assert stats["shape_count"] >= 2
        assert stats["total_chars"] > 0


def test_replace_text():
    processor = PptxProcessor()
    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "src.pptx")
        out = os.path.join(tmpdir, "out.pptx")
        _create_simple_pptx(src, [["原始标题", "原始内容", "保持不变"]])

        # slide 0, shape 0, para 0 和 para 1
        processor.replace_text(src, {"0:0:0": "新标题", "0:0:1": "新内容"}, out)

        result = processor.parse(out)
        paras = result["slides"][0]["shapes"][0]["paragraphs"]
        assert paras[0]["text"] == "新标题"
        assert paras[1]["text"] == "新内容"
        assert paras[2]["text"] == "保持不变"


def test_apply_format_changes():
    processor = PptxProcessor()
    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "src.pptx")
        out = os.path.join(tmpdir, "out.pptx")
        _create_simple_pptx(src, [["测试段落"]])

        changes = [
            {"target": "font_size", "scope": "all", "value_pt": 18},
            {"target": "line_spacing", "scope": "all", "value": 1.5},
        ]
        processor.apply_format_changes(src, changes, out)

        assert os.path.exists(out)
        prs = Presentation(out)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.size == Pt(18)
