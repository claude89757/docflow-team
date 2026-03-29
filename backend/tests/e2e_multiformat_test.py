"""E2E 多格式处理器集成测试

验证 get_processor → parse → replace_text → apply_format_changes 完整链路。
不依赖 Agent Teams / API key，纯本地处理器验证。
"""

import os
import tempfile

from docx import Document

from backend.processors import get_processor


def test_docx_full_flow():
    """docx: parse → replace → format → 验证输出"""
    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "input.docx")
        edited = os.path.join(tmpdir, "edited.docx")
        formatted = os.path.join(tmpdir, "formatted.docx")

        # 创建测试文档
        doc = Document()
        doc.add_heading("AI 生成的报告", level=1)
        doc.add_paragraph("此外，值得注意的是，人工智能技术在各个领域得到了广泛应用。")
        doc.add_paragraph("综上所述，我们可以得出结论。")
        doc.save(src)

        # 处理链路
        proc = get_processor(src)
        parsed = proc.parse(src)
        assert parsed["format"] == "docx"
        assert parsed["paragraph_count"] >= 2

        proc.replace_text(
            src, {"1": "AI 技术已经在医疗、教育等领域落地。", "2": "总结：技术落地需要结合场景。"}, edited
        )
        proc.apply_format_changes(edited, [{"target": "line_spacing", "scope": "all", "value": 1.15}], formatted)

        # 验证
        result_doc = Document(formatted)
        texts = [p.text for p in result_doc.paragraphs if p.text.strip()]
        assert any("AI 技术" in t for t in texts)
        assert any("总结" in t for t in texts)
        print(f"✓ docx flow: {len(texts)} paragraphs, output at {formatted}")


def test_pptx_full_flow():
    """pptx: parse → replace → format → 验证输出"""
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "input.pptx")
        edited = os.path.join(tmpdir, "edited.pptx")
        formatted = os.path.join(tmpdir, "formatted.pptx")

        # 创建测试 pptx
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = tx_box.text_frame
        tf.text = "此外，值得注意的是"
        p = tf.add_paragraph()
        p.text = "综上所述，我们认为"
        prs.save(src)

        # 处理链路
        proc = get_processor(src)
        parsed = proc.parse(src)
        assert parsed["format"] == "pptx"
        assert parsed["slide_count"] == 1

        proc.replace_text(src, {"0:0:0": "关键发现", "0:0:1": "我们的结论是"}, edited)
        proc.apply_format_changes(edited, [{"target": "font_size", "scope": "all", "value_pt": 16}], formatted)

        # 验证
        result_prs = Presentation(formatted)
        slide = result_prs.slides[0]
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text
        assert "关键发现" in all_text
        print(f"✓ pptx flow: {parsed['slide_count']} slides, output at {formatted}")


def test_xlsx_full_flow():
    """xlsx: parse → replace (跳过公式) → format → 验证输出"""
    from openpyxl import Workbook, load_workbook

    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "input.xlsx")
        edited = os.path.join(tmpdir, "edited.xlsx")
        formatted = os.path.join(tmpdir, "formatted.xlsx")

        # 创建测试 xlsx
        wb = Workbook()
        ws = wb.active
        ws.title = "数据"
        ws["A1"] = "此外，报告标题"
        ws["B1"] = "描述文字"
        ws["A2"] = 100
        ws["B2"] = "=A2*2"  # 公式
        wb.save(src)
        wb.close()

        # 处理链路
        proc = get_processor(src)
        parsed = proc.parse(src)
        assert parsed["format"] == "xlsx"
        assert parsed["sheet_count"] == 1

        proc.replace_text(src, {"数据:1:1": "季度报告", "数据:1:2": "Q1 数据汇总"}, edited)
        proc.apply_format_changes(edited, [{"target": "font_size", "scope": "all", "value": 12}], formatted)

        # 验证
        wb = load_workbook(formatted)
        ws = wb.active
        assert ws["A1"].value == "季度报告"
        assert ws["B1"].value == "Q1 数据汇总"
        assert ws["A2"].value == 100  # 数字不动
        assert ws["B2"].value == "=A2*2"  # 公式不动
        wb.close()
        print(f"✓ xlsx flow: {parsed['sheet_count']} sheets, formulas preserved, output at {formatted}")


def test_pdf_full_flow():
    """pdf: parse → to_docx → replace → format → 验证输出为 docx"""
    import fitz

    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "input.pdf")
        edited = os.path.join(tmpdir, "edited.pdf")
        formatted = os.path.join(tmpdir, "formatted.pdf")

        # 创建测试 PDF
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Report Title", fontsize=18)
        page.insert_text((72, 100), "This is the body text of the document.", fontsize=11)
        page.insert_text((72, 120), "Another paragraph with more content.", fontsize=11)
        doc.save(src)
        doc.close()

        # 处理链路
        proc = get_processor(src)
        parsed = proc.parse(src)
        assert parsed["format"] == "pdf"
        assert parsed["page_count"] == 1

        # replace_text 会转 docx
        result = proc.replace_text(src, {"0": "Updated Title"}, edited)
        assert result.endswith(".docx")

        # apply_format 也会转 docx
        result2 = proc.apply_format_changes(src, [{"target": "line_spacing", "scope": "all", "value": 1.2}], formatted)
        assert result2.endswith(".docx")

        # 验证 to_docx 独立功能
        docx_path = os.path.join(tmpdir, "converted.docx")
        proc.to_docx(src, docx_path)
        from docx import Document as DocxDocument

        converted = DocxDocument(docx_path)
        all_text = " ".join(p.text for p in converted.paragraphs)
        assert "Report Title" in all_text
        print(f"✓ pdf flow: {parsed['page_count']} pages, output as docx at {result}")


def test_processor_factory_routing():
    """验证工厂函数正确路由到各处理器"""
    from backend.processors.docx_processor import DocxProcessor
    from backend.processors.pdf_processor import PdfProcessor
    from backend.processors.pptx_processor import PptxProcessor
    from backend.processors.xlsx_processor import XlsxProcessor

    assert isinstance(get_processor("test.docx"), DocxProcessor)
    assert isinstance(get_processor("test.DOCX"), DocxProcessor)
    assert isinstance(get_processor("/path/to/file.pptx"), PptxProcessor)
    assert isinstance(get_processor("data.xlsx"), XlsxProcessor)
    assert isinstance(get_processor("report.pdf"), PdfProcessor)

    try:
        get_processor("file.txt")
        raise AssertionError("Should raise ValueError")
    except ValueError:
        pass
    print("✓ factory routing: all formats correctly dispatched")


if __name__ == "__main__":
    test_processor_factory_routing()
    test_docx_full_flow()
    test_pptx_full_flow()
    test_xlsx_full_flow()
    test_pdf_full_flow()
    print("\n=== All E2E multi-format tests passed ===")
